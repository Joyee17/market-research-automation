"""
市场调研自动化工具 — 完整版 Python 脚本
=====================================================
功能：
  1. 读取内部数据输入模板（用户填写的信息）
  2. 自动判断调研类型（A/B/C/D/E/F）
  3. 针对每个模块自动搜索权威外部数据
  4. 调用 Claude API 生成 McKinsey 级别分析
  5. 自动生成 Word 报告（.docx）
  6. 自动生成 PowerPoint 汇报（.pptx）

依赖安装：
  pip install anthropic python-docx python-pptx requests

运行方式：
  python market_research_tool.py

作者：市场调研自动化系统 v1.0
"""

import os
import json
import time
import re
from datetime import datetime
from pathlib import Path
import anthropic
from docx import Document
from docx.shared import Pt, RGBColor, Inches, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor as PPTXColor
from pptx.enum.text import PP_ALIGN


# ─────────────────────────────────────────────────────────────
# CONFIG：每次调研只需要修改这里
# ─────────────────────────────────────────────────────────────

CONFIG = {
    # API Key（也可设置环境变量 ANTHROPIC_API_KEY）
    "api_key": os.environ.get("ANTHROPIC_API_KEY", "your-api-key-here"),

    # 输入文件路径（填写好的内部数据模板）
    "input_file": "00_内部数据输入模板.md",

    # 输出目录
    "output_dir": "output_reports",

    # 模型选择
    "model": "claude-sonnet-4-20250514",

    # 每模块最大token
    "max_tokens": 1500,

    # 是否开启网络搜索增强（需要工具支持）
    "enable_web_search": True,

    # 报告语言
    "language": "中文",
}

# ─────────────────────────────────────────────────────────────
# 权威数据来源层级表（嵌入提示词，确保数据质量）
# ─────────────────────────────────────────────────────────────

AUTHORITY_SOURCES = """
## 数据来源权威性要求（严格遵守）

### L1 级（最高权威，优先引用）
- 国际机构：World Bank, IMF, IEA, IRENA, WTO, UN, OECD
- 政府官方：各国政府统计局、能源部、财政部
- 顶级咨询：McKinsey Global Institute, BCG, Bain, Deloitte, PwC, EY
- 金融监管：美联储, ECB, BIS

### L2 级（高权威，常规引用）
- 顶级行业研究：BloombergNEF, Gartner, Forrester, IDC, Wood Mackenzie
- 权威媒体：路透社, 彭博社, 《金融时报》, 《经济学人》, WSJ
- 行业协会：SEIA, BNEF, 各国光伏协会, 矿业协会
- 学术期刊：Nature Energy, Science, PNAS（同行评审）

### L3 级（中等权威，须交叉验证）
- 企业年报（上市公司）、可持续发展报告
- 权威行业媒体：PV Tech, Renew Economy, Mining Technology
- 平台数据：Global Solar Atlas (World Bank支持), NASA POWER

### L4 级（参考级，必须标注[待验证]）
- Statista, IBISWorld, Grand View Research（市场规模数据）
- LinkedIn, Crunchbase（企业信息）
- 非同行评审报告

### 禁止引用
- 无来源数据、AI生成内容、超过24个月的市场规模数据
- 单一来源的市场份额数据（必须交叉验证）
"""

# ─────────────────────────────────────────────────────────────
# 调研类型判断逻辑
# ─────────────────────────────────────────────────────────────

TYPE_MODULES = {
    "A": {
        "name": "新市场进入研究（Market Entry Research）",
        "modules": [
            ("exec_summary", "执行摘要"),
            ("product_anchor", "产品定位锚定"),
            ("market_screening", "市场筛选评分矩阵"),
            ("market_deep_dive", "重点市场深度分析"),
            ("icp", "目标客户画像（ICP）"),
            ("competitive_winrate", "竞争格局与差异化胜率"),
            ("business_model", "商业模式设计"),
            ("gtm_plan", "市场进入路径（GTM）"),
            ("decision_matrix", "优先级决策矩阵与最终结论"),
        ]
    },
    "B": {
        "name": "竞争情报研究（Competitive Intelligence）",
        "modules": [
            ("exec_summary", "执行摘要"),
            ("competitive_landscape", "竞争格局总览"),
            ("competitor_profile", "竞品深度画像"),
            ("winrate_analysis", "差异化胜率分析"),
            ("battlecard", "销售Battlecard"),
            ("counter_plan", "反制行动计划"),
        ]
    },
    "C": {
        "name": "客户需求研究（Customer & Demand Insight）",
        "modules": [
            ("exec_summary", "执行摘要"),
            ("research_design", "研究设计与访谈框架"),
            ("demand_analysis", "需求优先级矩阵"),
            ("icp_update", "客户画像精确化"),
            ("willingness_to_pay", "支付意愿与定价洞察"),
            ("gtm_message", "GTM信息优化建议"),
        ]
    },
    "D": {
        "name": "市场深耕研究（Market Penetration）",
        "modules": [
            ("exec_summary", "执行摘要"),
            ("growth_diagnosis", "增长结构诊断"),
            ("customer_health", "客户健康度分析"),
            ("untapped_market", "渗透率空白识别"),
            ("growth_levers", "增长杠杆ROI排序"),
            ("action_plan", "季度行动计划"),
        ]
    },
    "E": {
        "name": "产品与定价研究（Product & Pricing）",
        "modules": [
            ("exec_summary", "执行摘要"),
            ("product_assessment", "产品功能价值评估"),
            ("roadmap_priority", "产品路线图优先级矩阵"),
            ("pricing_diagnosis", "定价现状诊断"),
            ("willingness_to_pay_e", "支付意愿研究"),
            ("pricing_structure", "定价结构与实施路径"),
        ]
    },
    "F": {
        "name": "品牌与传播研究（Brand & Communication）",
        "modules": [
            ("exec_summary", "执行摘要"),
            ("brand_awareness", "品牌认知漏斗分析"),
            ("brand_association", "品牌联想与竞品对比"),
            ("message_effectiveness", "核心信息有效性"),
            ("channel_efficiency", "渠道触达效率分析"),
            ("brand_strategy", "品牌战略与90天行动"),
        ]
    }
}


def detect_research_type(content: str) -> str:
    """根据输入内容自动判断调研类型"""
    content_lower = content.lower()

    # 明确选择
    for t in ["A", "B", "C", "D", "E", "F"]:
        if f"[x] {t}型" in content or f"[x]{t}型" in content:
            return t

    # 关键词判断
    if any(w in content for w in ["新市场", "进入", "market entry", "候选市场"]):
        return "A"
    if any(w in content for w in ["竞品", "竞争情报", "battlecard", "competitive"]):
        return "B"
    if any(w in content for w in ["客户需求", "用户访谈", "支付意愿", "demand"]):
        return "C"
    if any(w in content for w in ["市场深耕", "增长", "流失率", "penetration"]):
        return "D"
    if any(w in content for w in ["产品定价", "路线图", "定价策略", "pricing"]):
        return "E"
    if any(w in content for w in ["品牌", "传播", "认知度", "brand"]):
        return "F"

    return "A"  # 默认新市场进入


def build_system_prompt(research_type: str) -> str:
    """构建系统提示词"""
    type_info = TYPE_MODULES[research_type]
    return f"""你是麦肯锡资深市场策略顾问，专精B2B工业/能源/技术产品的高管汇报级市场调研。

当前调研类型：{type_info['name']}

{AUTHORITY_SOURCES}

## 输出质量红线（违反则返工）

1. **结论必须明确**：禁止「各有优劣」「视情况而定」「需进一步研究」
2. **每个洞察完成四层**：What（观察）→ Why（原因）→ So What（含义）→ Now What（行动）
3. **数据必须有来源**：每个关键数据附来源机构+年份，L4级数据标注[待验证]
4. **客户描述ICP级别**：行业+规模+地理+场景，禁止泛化（如「数据中心企业」）
5. **竞品分析场景化**：在[具体条件]下，[竞品]问题[数据]，我们能做到[数据]，客户收益[可量化]
6. **行动必须可执行**：时间节点+具体动作+负责方+量化成功标准
7. **风险有应对**：预警信号+预防措施+止损触发点
8. **口径全文一致**：同一指标只用一个单位，报告开头建立口径说明

## 叙事结构
金字塔原则：结论先行 → 核心论点 → 支撑数据 → 具体行动
"""


def build_module_prompt(module_key: str, context: str, research_type: str) -> str:
    """为每个模块构建专属提示词"""

    PROMPTS = {
        "exec_summary": f"""基于以下调研背景，生成「执行摘要」。
要求：不超过500字，高管5分钟内读完可直接拍板，Markdown格式。

{context}

必须包含：
## 核心结论
推荐方向（1-2个）+ 3条数据支撑理由（每条必须引用权威数据来源）

## 次选与暂缓
各一句话，原因具体，附数据

## 第一步行动
[具体时间节点] → [做什么] → [找谁] → [预期可量化结果]

## 90天里程碑
3个可量化检查点，含负责方""",

        "product_anchor": f"""基于以下信息，生成「产品定位锚定」模块（整份报告的技术基准）。

{context}

必须包含：
## 核心技术指标对比表
（我们 vs 主要竞品，每个指标附权威数据来源，禁止无来源数据）

## 赢的场景（2-3个）
每个场景格式：
- 场景条件：[具体环境/客户/规模]
- 竞品在此场景的量化问题：[数据+来源]
- 我们的量化优势：[数据+来源]
- 转化为客户财务收益：[具体数字]

## 输的场景（1-2个，诚实写出）
对市场筛选的含义

## 权重建议
对后续市场评分矩阵，哪个维度应设最高权重，理由""",

        "market_screening": f"""基于以下信息，生成「市场筛选评分矩阵」。

{context}

六维度评分（1-5分）：需求强度25%、技术匹配25%、支付能力20%、进入难度15%、竞争强度10%、战略价值5%

要求：
- 每个分数必须有具体数据支撑（来源标注）
- 禁止凭印象打分

输出：
## 评分矩阵表格
（各市场×各维度，含加权总分，每格附一句话依据+数据来源）

## 市场分层结论
立即进入（>4.0）/观察跟进（3.0-4.0）/暂缓（<3.0）
每个分层：3条具体理由，每条附数据来源

## 唯一最优市场
如果只能做一个市场，推荐哪个，为什么（必须明确，附3条核心理由）

## 暂缓市场重启条件
什么条件下暂缓市场可以重新评估""",

        "market_deep_dive": f"""基于以下信息，对评分最高的市场生成「深度分析」。

{context}

必须包含：
## 市场本质（一句话锚点）
格式：「[市场] = [客户核心痛点] × [现有方案结构性缺陷] × [我们的独特价值] 的交叉市场」

## 供需结构与三大核心矛盾
每个矛盾：现状（数据+L1/L2来源）→ 本质问题 → 对我们的具体机会

## 市场规模（TAM→SAM→SOM）
- TAM：总规模，L1/L2来源，附推导逻辑
- SAM：产品适用条件筛选后，附筛选比例依据
- SOM：接触量×转化率×客户价值，每个假设明确标注

## 标杆案例（1-2个）
成功因素+局限性（有数据）+我们的超越机会

## 风险矩阵
概率×影响→预警信号→预防→止损触发点""",

        "icp": f"""基于以下信息，生成「目标客户画像（ICP）」。
定义：最可能第一批成交且能产生示范效应的客户类型。

{context}

必须包含：
## 客户特征（ICP级别，禁止泛化）
行业+子行业+规模量化指标+具体地理位置（城市级）+组织特征

## 核心痛点（3个）
格式：现状（客户原话/公开资料引用）→ 量化损失 → 我们能解决的程度

## 购买触发器（2-3个具体事件场景）

## 决策链表格
角色/典型职位/决策作用/核心关注点/最担心的风险

## 供应商评估标准（按重要性排序）

## 典型代表企业（3-5家，附符合ICP的具体依据）

## 触达路径
渠道+时机+切入话题（禁止写「介绍我们产品」）""",

        "competitive_winrate": f"""基于以下信息，生成「竞争格局与差异化胜率」分析。

{context}

必须包含：
## 竞争阵营分析表
竞争类型/代表企业/市场份额（来源）/核心弱点（场景化）/我们的非对称优势

## 三场景差异化胜率（严格格式）
场景①：在[具体条件]下，[竞品]的问题是[量化数据+来源]，
我们能做到[量化数据+来源]，对客户意味着[可量化财务收益]

## 应主动避开的场景（我们结构性劣势）

## 差异化一句话定位
「[我们]是[ICP]在[具体场景]下比[竞品]更好的[核心价值]方案」

## 非对称优势的可持续性
竞品复制我们优势需要多长时间，为什么""",

        "business_model": f"""基于以下信息，生成「商业模式设计」。

{context}

## 模式适用性对比表
（设备直销/EPC/PPA/EaaS/合资，含：适配客户/资金要求/客户门槛/首单难度/规模化潜力）

## 第一阶段推荐模式
必须明确推荐（禁止各有优劣），附：
- 推荐理由3条（每条有数据支撑）
- 关键成功要素2-3个
- 典型合同结构要点
- 预期首单规模区间

## 第二阶段演进方向（12-36个月）

## 必须提前准备的前提条件清单
（认证/合作伙伴/融资，每项附时间估算）""",

        "gtm_plan": f"""基于以下信息，生成「市场进入路径（GTM）」。

{context}

## 进入路径选择与推荐
（直接/合作伙伴/展会/标杆项目，推荐明确+理由+避开路径）

## 客户触达策略
- ICP在哪里获取行业信息（具体渠道名称）
- 我们如何建立第一次有效接触（具体方式）
- 最佳切入时机与话题

## 分阶段行动计划
第一阶段（0-6月）：核心目标+具体动作（做什么/找谁/方式/预期结果）+量化成功标准
第二阶段（6-18月）：同格式
第三阶段（18月+）：规模复制逻辑

## 启动前必须准备清单""",

        "decision_matrix": f"""基于以下信息，生成「优先级决策矩阵与最终结论」（报告最重要的一页）。

{context}

## 市场优先级决策矩阵
表格：机会评级/进入时机/推荐模式/关键前提条件/风险等级

## 最终战略结论（必须明确，绝不模糊）
第一优先：[市场]，理由①②③（每条含权威数据来源）
第二梯队：[市场]，进入条件
暂缓：[市场]，具体原因（一句话）

## 第一步行动（最终落地）
[时间] [具体动作] [找谁/通过什么渠道] [预期可量化结果]

## 90天里程碑
时间点/具体里程碑/量化成功标准/负责方""",

        # B型模块
        "competitive_landscape": f"""基于以下信息，生成「竞争格局总览」（B型竞争情报）。

{context}

## 竞争阵营三层分类
直接竞品/间接竞品/潜在进入者，含：市场份额（来源）+威胁等级+监控优先级

## 市场份额动态（过去12个月）
变化数据+驱动因素+我们相对地位趋势（均需来源）

## 最高威胁竞品定性
具体威胁场景（非笼统「综合竞争力强」）

## 竞争格局未来12个月预判（更激烈/分散，驱动因素）

## 资源分配决策含义""",

        "competitor_profile": f"""基于以下信息，生成「竞品深度画像」（B型）。

{context}

对每个主要竞品：
## 战略意图分析
融资/招聘方向/新市场信号→战略推断→置信度H/M/L

## 产品能力场景化对比（禁止纯参数表）
每个对比绑定具体使用场景，附第三方来源

## 四类弱点深挖
客户公开抱怨（引用原文）/技术局限（有数据）/服务弱点/覆盖空白

## 未来6个月最可能重大动作预判""",

        "winrate_analysis": f"""基于以下信息，生成「差异化胜率分析」（B型）。

{context}

## 竞争定位地图
按ICP最在意的2个维度（非技术维度）定位各竞品

## 三场景差异化胜率
严格格式：在[条件]下，[竞品]问题[数据+来源]，我们能做到[数据+来源]，客户收益[可量化]

## 应避开的竞争场景（结构性弱点）

## 差异化优势可持续性评估""",

        "battlecard": f"""基于以下信息，生成「销售Battlecard」（B型，单页浓缩）。

{context}

## 一句话定位差异

## 我们的3个核心优势（每个附数据+客户可量化价值）

## 竞品声称 vs 我们的应对话术 vs 支撑证据（表格）

## 客户常见异议 vs 专业回应 vs 关键证据（表格）

## 主动推的场景 vs 需要转移的场景""",

        "counter_plan": f"""基于以下信息，生成「反制行动计划」（B型）。

{context}

## 竞品近期重大动作的即时应对
按H/M/L威胁等级，含时间节点和负责方

## 主动进攻计划
从竞品手中赢得客户：识别标准+切入时机+切入话术

## CI持续监控机制
月度监控清单+季度深度更新触发条件""",

        # C型模块
        "research_design": f"""基于以下信息，生成「研究设计与访谈框架」（C型）。

{context}

## 访谈对象分层设计
现有客户/流失客户/竞品客户/非客户，各类型数量+找到方式

## 完整访谈问题框架
现状探索/痛点深挖/解决方案评估/Van Westendorp四问法

## 三层需求识别框架
表达需求/真实需求/潜在需求的区别与识别方法""",

        "demand_analysis": f"""基于以下信息，生成「需求优先级矩阵」（C型）。

{context}

## 需求优先级评分矩阵
提及频率×痛感强度×支付意愿×当前满足度（1-5分，综合优先级计算）

## 四象限分类
立即解决/保持优势/观察跟进/过度投资

## 最被低估的需求（高优先级但当前未解决）

## 不值得追逐的需求（表达强烈但支付意愿低）

## 产品迭代决策含义""",

        "icp_update": f"""基于以下信息，生成「客户画像精确化」（C型）。

{context}

## ICP更新对比表
原有假设 vs 研究发现 vs 更新后定义

## 高价值 vs 低价值客户分化特征

## 需要新增的ICP细分（研究发现的新客户群）

## 需要排除的客户画像（与能力不匹配）

## 销售资质审核标准更新建议""",

        "willingness_to_pay": f"""基于以下信息，生成「支付意愿与定价洞察」（C型）。

{context}

## Van Westendorp四价格点
拒绝价/可接受上限/便宜购买/怀疑价+建议定价区间

## 分客户群支付意愿差异

## 当前定价的位置诊断（是否存在价值低估）

## 溢价空间分析（为哪些具体特性，幅度多大）

## 定价调整建议与风险""",

        "gtm_message": f"""基于以下信息，生成「GTM信息优化建议」（C型）。

{context}

## 客户语言 vs 我们的语言对比表
当前表达 vs 客户实际用语 vs 调整建议

## 核心价值主张重构（基于客户语言）

## 分决策链角色的差异化信息

## 最有力的3个客户引言（可直接用于营销）""",

        # D型模块
        "growth_diagnosis": f"""基于以下信息，生成「增长结构诊断」（D型）。

{context}

## 增长来源分解
新客获取/现有客户扩展/价格调整/市场自然增长（各占比+趋势+数据来源）

## 增长漏出分析
客户流失/降级/未转化线索（各占比+根本原因）

## 净增长真实来源（「我们做好了」vs「市场在涨」的分离）

## 市场份额动态（相对竞品的地位，非绝对数字）

## 增长质量评估与战略含义""",

        "customer_health": f"""基于以下信息，生成「客户健康度分析」（D型）。

{context}

## 客户价值分层（A/B/C/D级定义+各层客户数+ARR占比+策略）

## 健康度评分模型（使用频率/续约意愿/NPS/扩展行为/服务质量，权重）

## 流失风险前10名（干预计划）

## 扩展机会前10名（最容易的增收来源）

## 流失根因分析（可阻止 vs 不可阻止+比例+应对）""",

        "untapped_market": f"""基于以下信息，生成「渗透率空白识别」（D型）。

{context}

## 细分市场渗透率矩阵
按行业/规模/地区/场景，各细分潜在客户数+当前覆盖+渗透率+未覆盖原因

## 最大渗透率空白及进入障碍分析

## 「一个改变解锁多个客户群」的机会

## ROI最高的空白攻克路径""",

        "growth_levers": f"""基于以下信息，生成「增长杠杆ROI排序」（D型）。

{context}

## 增长杠杆全景（含预期贡献+所需投入+ROI）

## 每个杠杆四层分析
What量化→Why为什么有效→So What不做的损失→Now What第一步

## 最终优先级决策表（按ROI，含启动时间和负责方）

## 资源分配建议（防御%/进攻%/探索%+理由）""",

        "action_plan": f"""基于以下信息，生成「季度行动计划」（D型）。

{context}

## 90天执行路线图
第1/2/3个月：核心目标+具体动作+量化成功标准

## 关键里程碑表（成功标准/时间/负责方）

## 风险预案（每个主要行动的风险+应急措施）""",

        # E型模块
        "product_assessment": f"""基于以下信息，生成「产品功能价值评估」（E型）。

{context}

## Kano模型功能分类
必备/期望/魅力/无差异（当前满足度+竞品对比+数据来源）

## 被错误投入资源的功能（无差异功能清单）

## 被低估的魅力功能（未宣传的差异化资产）

## 产品使用数据分析
核心功能使用率/Aha Moment/流失预测信号/Time-to-Value

## 迭代资源分配建议""",

        "roadmap_priority": f"""基于以下信息，生成「产品路线图优先级矩阵」（E型）。

{context}

## 需求来源汇总（客户/败单/流失/竞品差距，真实性验证）

## 优先级评分矩阵
市场需求强度×收入影响×战略价值×实现复杂度×竞争紧迫性（权重×评分）

## 四象限决策（立即做/计划做/再议/不做）

## 「不做」决策及明确理由

## Quick Win识别（2-4周高价值改善）

## 路线图与定价联动逻辑""",

        "pricing_diagnosis": f"""基于以下信息，生成「定价现状诊断」（E型）。

{context}

## 价值-价格对齐分析
每个功能的客户感知价值 vs 是否计费 → 价值捕获状态

## 最大价值捕获缺口（免费提供的高价值功能）

## 过度定价风险（单独计费但感知低的功能）

## 竞品定价基准对比（我们的位置，是否与定位一致）

## 定价溢价的合理空间""",

        "willingness_to_pay_e": f"""基于以下信息，生成「支付意愿研究」（E型）。

{context}

## Van Westendorp四价格点（基于真实访谈数据，n≥15）

## 分客户层支付意愿差异

## 当前定价位置诊断

## 价格弹性评估（涨价10/20/30%的预估流失率）

## 定价分层机会""",

        "pricing_structure": f"""基于以下信息，生成「定价结构与实施路径」（E型）。

{context}

## 定价结构选项评估
按席位/用量/功能分层/成果分成/混合（适配性分析）

## 最优定价结构推荐（必须明确，附与客户价值感知的对齐逻辑）

## 定价分层设计（分割点依据）

## 涨价实施路径（价值强化→新客试行→现有客通知→全面执行→效果评估）

## 风险缓冲方案""",

        # F型模块
        "brand_awareness": f"""基于以下信息，生成「品牌认知漏斗分析」（F型）。

{context}

## 品牌认知五层漏斗
知名度/理解度/差异化认知/可信度/偏好度（当前数据+目标+差距+测量方法）

## 最大漏斗断层识别与商业影响量化

## 与竞品的对比（竞品在哪层领先）

## 品牌投入方向决策含义""",

        "brand_association": f"""基于以下信息，生成「品牌联想与竞品对比感知」（F型）。

{context}

## 品牌联想词云分析
客户实际高频词 vs 我们希望的 vs 差距

## 「品牌现实」vs「品牌意图」最大差距

## 意外的正面联想（被忽视的差异化资产）

## 需要纠正的负面联想

## 竞品品牌对比感知矩阵""",

        "message_effectiveness": f"""基于以下信息，生成「核心信息有效性」（F型）。

{context}

## 当前核心信息测试结果（理解率/共鸣率/可信度/行动意愿）

## 高共鸣 vs 无效信息分类（附客户原话）

## 「最有力的一句话」

## 「技术语言陷阱」识别

## 核心价值主张更新建议（分角色）""",

        "channel_efficiency": f"""基于以下信息，生成「渠道触达效率分析」（F型）。

{context}

## 目标受众信息接触点调研（使用频率/信任度/我们存在感/投入产出）

## 最高价值渠道（高信任度+我们低存在感）

## 资源错配渠道（高投入但低受众覆盖）

## 「品牌沙漠」识别

## 下年度营销预算分配建议""",

        "brand_strategy": f"""基于以下信息，生成「品牌战略与90天行动」（F型）。

{context}

## 品牌定位声明（内部格式）
对于[ICP]，面对[痛点]，[品牌]是[类别]，提供[差异化价值]，因为[独特能力]，不同于[竞品]

## 三层信息体系（核心主张/分层信息/证明层）

## 90天行动计划（每步有量化成功标准）

## 品牌投资ROI预期""",
    }

    base_prompt = f"""
调研背景信息：
{context}

{PROMPTS.get(module_key, f'生成「{module_key}」模块，严格遵守麦肯锡高管汇报质量标准。')}
"""
    return base_prompt


# ─────────────────────────────────────────────────────────────
# Claude API 调用
# ─────────────────────────────────────────────────────────────

def call_claude(client: anthropic.Anthropic, prompt: str, system: str, cfg: dict) -> str:
    """调用 Claude API，支持网络搜索"""
    tools = []
    if cfg.get("enable_web_search"):
        tools = [{"type": "web_search_20250305", "name": "web_search"}]

    try:
        if tools:
            msg = client.messages.create(
                model=cfg["model"],
                max_tokens=cfg["max_tokens"],
                system=system,
                tools=tools,
                messages=[{"role": "user", "content": prompt}],
            )
        else:
            msg = client.messages.create(
                model=cfg["model"],
                max_tokens=cfg["max_tokens"],
                system=system,
                messages=[{"role": "user", "content": prompt}],
            )

        # 提取文本内容（包含工具调用后的最终回复）
        text_parts = [b.text for b in msg.content if hasattr(b, "text") and b.text]
        return "\n".join(text_parts)

    except Exception as e:
        return f"[生成失败] {str(e)}"


# ─────────────────────────────────────────────────────────────
# Word 文档生成
# ─────────────────────────────────────────────────────────────

def set_cell_background(cell, hex_color):
    """设置表格单元格背景色"""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'), 'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'), hex_color)
    tcPr.append(shd)


def add_styled_heading(doc, text, level, color_hex="1B3A6B"):
    """添加带颜色的标题"""
    heading = doc.add_heading(text, level=level)
    for run in heading.runs:
        run.font.color.rgb = RGBColor(
            int(color_hex[0:2], 16),
            int(color_hex[2:4], 16),
            int(color_hex[4:6], 16)
        )
    return heading


def generate_word_report(
    modules_output: dict,
    module_names: dict,
    research_type: str,
    context_summary: dict,
    output_path: str
):
    """生成专业 Word 报告"""
    doc = Document()

    # 页面设置
    from docx.shared import Cm
    section = doc.sections[0]
    section.page_width = Cm(21)
    section.page_height = Cm(29.7)
    section.left_margin = Cm(2.5)
    section.right_margin = Cm(2.5)
    section.top_margin = Cm(2.5)
    section.bottom_margin = Cm(2.5)

    # ── 封面页 ──
    doc.add_paragraph()
    doc.add_paragraph()

    title_para = doc.add_paragraph()
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_run = title_para.add_run("市场调研报告")
    title_run.font.size = Pt(28)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(27, 58, 107)

    doc.add_paragraph()

    type_name = TYPE_MODULES[research_type]["name"]
    subtitle_para = doc.add_paragraph()
    subtitle_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle_run = subtitle_para.add_run(type_name)
    subtitle_run.font.size = Pt(16)
    subtitle_run.font.color.rgb = RGBColor(70, 130, 180)

    doc.add_paragraph()
    doc.add_paragraph()

    # 基本信息表
    info_table = doc.add_table(rows=5, cols=2)
    info_table.style = "Table Grid"
    info_data = [
        ("产品/服务", context_summary.get("product", "—")),
        ("目标市场", context_summary.get("markets", "—")),
        ("调研类型", f"{research_type}型：{type_name}"),
        ("生成日期", datetime.now().strftime("%Y年%m月%d日")),
        ("质量标准", "McKinsey 高管汇报级别"),
    ]
    for i, (label, value) in enumerate(info_data):
        row = info_table.rows[i]
        row.cells[0].text = label
        row.cells[1].text = value
        row.cells[0].paragraphs[0].runs[0].font.bold = True
        set_cell_background(row.cells[0], "D6E4F0")

    doc.add_page_break()

    # ── 目录提示 ──
    doc.add_heading("目录", level=1)
    for key, name in module_names.items():
        if key in modules_output:
            p = doc.add_paragraph(style="List Bullet")
            p.add_run(name)

    doc.add_page_break()

    # ── 各模块内容 ──
    for key, name in module_names.items():
        if key not in modules_output:
            continue

        content = modules_output[key]

        # 模块标题
        doc.add_heading(name, level=1)

        # 将 Markdown 内容转换为 Word 格式
        lines = content.split("\n")
        i = 0
        while i < len(lines):
            line = lines[i].strip()

            if not line:
                i += 1
                continue

            if line.startswith("## "):
                doc.add_heading(line[3:], level=2)
            elif line.startswith("### "):
                doc.add_heading(line[4:], level=3)
            elif line.startswith("#### "):
                doc.add_heading(line[5:], level=4)
            elif line.startswith("| "):
                # 表格处理
                table_lines = []
                while i < len(lines) and lines[i].strip().startswith("|"):
                    table_lines.append(lines[i].strip())
                    i += 1

                if table_lines:
                    # 解析表格
                    rows = []
                    for tl in table_lines:
                        if "---" in tl:
                            continue
                        cells = [c.strip() for c in tl.split("|")[1:-1]]
                        if cells:
                            rows.append(cells)

                    if rows:
                        max_cols = max(len(r) for r in rows)
                        table = doc.add_table(rows=len(rows), cols=max_cols)
                        table.style = "Table Grid"
                        for ri, row_data in enumerate(rows):
                            for ci, cell_text in enumerate(row_data):
                                if ci < max_cols:
                                    table.rows[ri].cells[ci].text = cell_text
                                    if ri == 0:
                                        for run in table.rows[ri].cells[ci].paragraphs[0].runs:
                                            run.font.bold = True
                                        set_cell_background(table.rows[ri].cells[ci], "BDD7EE")
                        doc.add_paragraph()
                continue

            elif line.startswith("- ") or line.startswith("* "):
                p = doc.add_paragraph(style="List Bullet")
                p.add_run(line[2:])
            elif line.startswith("**") and line.endswith("**"):
                p = doc.add_paragraph()
                run = p.add_run(line[2:-2])
                run.font.bold = True
            else:
                doc.add_paragraph(line)

            i += 1

        doc.add_page_break()

    # ── 数据质量声明 ──
    doc.add_heading("数据质量声明", level=1)
    declaration = doc.add_paragraph()
    declaration.add_run(
        "本报告数据来源遵循权威性层级标准：\n"
        "• L1级（最高）：世界银行、IMF、IEA、IRENA、麦肯锡全球研究院等国际权威机构\n"
        "• L2级（高）：BloombergNEF、Gartner、路透社、彭博社等顶级行业研究机构\n"
        "• L3级（中）：企业年报、权威行业媒体\n"
        "• [待验证]标注：需要进一步交叉验证的数据\n\n"
        f"报告生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M')}\n"
        "报告标准：McKinsey 高管汇报级别"
    )

    doc.save(output_path)
    print(f"✅ Word 报告已生成：{output_path}")


# ─────────────────────────────────────────────────────────────
# PowerPoint 生成
# ─────────────────────────────────────────────────────────────

def generate_pptx_report(
    modules_output: dict,
    module_names: dict,
    research_type: str,
    context_summary: dict,
    output_path: str
):
    """生成专业 PowerPoint 汇报"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # 颜色主题（深蓝商务风）
    DARK_BLUE = PPTXColor(27, 58, 107)
    MED_BLUE = PPTXColor(70, 130, 180)
    LIGHT_BLUE = PPTXColor(214, 228, 240)
    WHITE = PPTXColor(255, 255, 255)
    DARK_GRAY = PPTXColor(50, 50, 50)
    ACCENT = PPTXColor(255, 140, 0)

    def add_text_box(slide, text, left, top, width, height,
                     font_size=12, bold=False, color=None, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
        return txBox

    def add_rect(slide, left, top, width, height, fill_color, line=False):
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if not line:
            shape.line.fill.background()
        return shape

    # ── 封面幻灯片 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白

    # 深蓝背景
    add_rect(slide, 0, 0, 13.33, 7.5, DARK_BLUE)
    # 装饰条
    add_rect(slide, 0, 5.5, 13.33, 0.08, ACCENT)

    # 标题
    add_text_box(slide, "市场调研报告",
                 1, 1.5, 11, 1.5,
                 font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    type_name = TYPE_MODULES[research_type]["name"]
    add_text_box(slide, type_name,
                 1, 3.2, 11, 0.8,
                 font_size=20, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

    # 基本信息
    info_text = (
        f"产品：{context_summary.get('product', '—')}   "
        f"市场：{context_summary.get('markets', '—')}   "
        f"日期：{datetime.now().strftime('%Y.%m.%d')}"
    )
    add_text_box(slide, info_text,
                 1, 4.2, 11, 0.6,
                 font_size=12, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide, "McKinsey 高管汇报级别 | 数据来源：L1-L2 权威机构",
                 1, 6.5, 11, 0.5,
                 font_size=10, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

    # ── 目录幻灯片 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.33, 1.2, DARK_BLUE)
    add_text_box(slide, "目录", 0.5, 0.2, 12, 0.8,
                 font_size=24, bold=True, color=WHITE)

    modules_list = [(k, v) for k, v in module_names.items() if k in modules_output]
    cols = 2
    per_col = (len(modules_list) + 1) // cols
    for idx, (key, name) in enumerate(modules_list):
        col = idx // per_col
        row = idx % per_col
        add_text_box(
            slide,
            f"{idx + 1:02d}. {name}",
            0.5 + col * 6.5,
            1.5 + row * 0.65,
            6.0, 0.55,
            font_size=13, color=DARK_GRAY
        )

    # ── 各模块幻灯片 ──
    for idx, (key, name) in enumerate(modules_list):
        content = modules_output[key]

        # 提取前500字作为幻灯片内容（PPT不能放太多文字）
        lines = [l.strip() for l in content.split("\n") if l.strip() and not l.startswith("#")]
        preview_lines = lines[:8]  # 取前8行要点

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 顶部色条
        add_rect(slide, 0, 0, 13.33, 1.3, DARK_BLUE)
        # 编号
        add_text_box(slide, f"{idx + 1:02d}", 0.3, 0.15, 0.8, 0.9,
                     font_size=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        # 模块名称
        add_text_box(slide, name, 1.2, 0.2, 11, 0.9,
                     font_size=22, bold=True, color=WHITE)

        # 内容区域
        content_text = "\n".join([f"• {l}" for l in preview_lines if l])
        add_text_box(slide, content_text,
                     0.5, 1.5, 12.3, 5.5,
                     font_size=13, color=DARK_GRAY)

        # 底部页码
        add_text_box(slide, f"{idx + 3} / {len(modules_list) + 2}",
                     11.5, 7.0, 1.5, 0.4,
                     font_size=10, color=DARK_GRAY, align=PP_ALIGN.RIGHT)

    # ── 结尾幻灯片 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.33, 7.5, DARK_BLUE)
    add_rect(slide, 0, 3.5, 13.33, 0.08, ACCENT)
    add_text_box(slide, "数据来源与质量声明",
                 1, 1.5, 11, 1, font_size=28, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(
        slide,
        "L1级：World Bank · IMF · IEA · IRENA · McKinsey Global Institute\n"
        "L2级：BloombergNEF · Gartner · Reuters · Bloomberg\n"
        "L3级：企业年报 · 权威行业媒体\n"
        "[待验证] 标注：需进一步交叉验证的数据",
        1, 2.8, 11, 2.5,
        font_size=14, color=LIGHT_BLUE, align=PP_ALIGN.CENTER
    )

    prs.save(output_path)
    print(f"✅ PPT 报告已生成：{output_path}")


# ─────────────────────────────────────────────────────────────
# 主流程
# ─────────────────────────────────────────────────────────────

def extract_context_summary(content: str) -> dict:
    """从输入模板中提取关键摘要信息"""
    summary = {}

    patterns = {
        "product": r"【产品/服务名称】\s*\n→\s*(.+)",
        "markets": r"【候选市场/目标市场】\s*\n→\s*(.+)",
        "decision": r"【决策问题】\s*\n→\s*(.+)",
        "constraints": r"【公司资源约束】\s*\n([\s\S]*?)(?=\n【)",
    }

    for key, pattern in patterns.items():
        match = re.search(pattern, content)
        if match:
            summary[key] = match.group(1).strip()

    return summary


def run_research(cfg: dict):
    """主运行函数"""
    print("\n" + "="*60)
    print("  市场调研自动化工具 v1.0")
    print("  McKinsey 高管汇报级别")
    print("="*60)

    # 读取输入文件
    input_path = Path(cfg["input_file"])
    if not input_path.exists():
        print(f"❌ 输入文件不存在：{cfg['input_file']}")
        print("请先填写 00_内部数据输入模板.md 后再运行")
        return

    with open(input_path, "r", encoding="utf-8") as f:
        input_content = f.read()

    print(f"✅ 已读取输入文件：{input_path}")

    # 提取摘要
    context_summary = extract_context_summary(input_content)

    # 判断调研类型
    research_type = detect_research_type(input_content)
    type_info = TYPE_MODULES[research_type]
    print(f"✅ 调研类型：{research_type}型 — {type_info['name']}")

    # 初始化 Claude 客户端
    client = anthropic.Anthropic(api_key=cfg["api_key"])

    # 构建系统提示词
    system_prompt = build_system_prompt(research_type)

    # 创建输出目录
    output_dir = Path(cfg["output_dir"])
    output_dir.mkdir(exist_ok=True)
    date_str = datetime.now().strftime("%Y%m%d_%H%M")
    product_tag = context_summary.get("product", "report")[:10].replace("/", "_").replace(" ", "_")

    # 逐模块生成
    modules_output = {}
    module_names = {k: v for k, v in type_info["modules"]}
    total = len(type_info["modules"])

    print(f"\n开始生成 {total} 个模块...\n")

    for i, (module_key, module_name) in enumerate(type_info["modules"], 1):
        print(f"[{i}/{total}] 正在生成「{module_name}」...", end="", flush=True)

        prompt = build_module_prompt(module_key, input_content, research_type)

        result = call_claude(client, prompt, system_prompt, cfg)
        modules_output[module_key] = result

        word_count = len(result)
        print(f" 完成（{word_count}字）")

        # 避免频率限制
        if i < total:
            time.sleep(1)

    print(f"\n✅ 所有模块生成完毕\n")

    # 生成 Word 报告
    word_path = str(output_dir / f"市场调研报告_{product_tag}_{date_str}.docx")
    generate_word_report(
        modules_output, module_names,
        research_type, context_summary,
        word_path
    )

    # 生成 PPT 报告
    pptx_path = str(output_dir / f"市场调研报告_{product_tag}_{date_str}.pptx")
    generate_pptx_report(
        modules_output, module_names,
        research_type, context_summary,
        pptx_path
    )

    # 保存 JSON 中间结果（便于调试和二次加工）
    json_path = str(output_dir / f"raw_output_{product_tag}_{date_str}.json")
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump({
            "research_type": research_type,
            "type_name": type_info["name"],
            "context_summary": context_summary,
            "modules": modules_output,
            "generated_at": datetime.now().isoformat(),
        }, f, ensure_ascii=False, indent=2)

    print(f"\n{'='*60}")
    print(f"✅ 全部完成！输出文件：")
    print(f"   📄 Word：{word_path}")
    print(f"   📊 PPT： {pptx_path}")
    print(f"   🗂  JSON：{json_path}")
    print(f"{'='*60}\n")


if __name__ == "__main__":
    run_research(CONFIG)
